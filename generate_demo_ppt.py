from pptx import Presentation
from pptx.util import Pt

# Create presentation
prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "AI Chatbot MVP – Demo & Edge Case Handling"
subtitle.text = "Demo of User Input → AI Response Flow and Edge Case Handling\nYour Name / Team\nDate"

# Slide 2: User Input → AI Response Flow
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "User Input → AI Response Flow"
content = (
    "• Step 1: User types a question\n"
    "    (Insert screenshot: chat input with message being typed)\n"
    "• Step 2: User's message appears, loading spinner is visible\n"
    "    (Insert screenshot: spinner/loading state)\n"
    "• Step 3: AI's response appears\n"
    "    (Insert screenshot: AI reply)"
)
slide.placeholders[1].text = content

# Slide 3: Handling Slow AI Response
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Handling Slow AI Response"
content = (
    "• Step 1: User submits a question\n"
    "    (Insert screenshot: user message)\n"
    "• Step 2: Spinner remains visible for several seconds\n"
    "    (Insert screenshot: spinner/loading state, possibly with a toast: 'Still working on your answer...')\n"
    "• Step 3: AI response appears after delay\n"
    "    (Insert screenshot: AI reply)"
)
slide.placeholders[1].text = content

# Slide 4: Error Handling
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Error Handling"
content = (
    "• Step 1: User submits a question\n"
    "    (Insert screenshot: user message)\n"
    "• Step 2: Backend/API error occurs\n"
    "    (Insert screenshot: error toast or message, e.g., 'Sorry, something went wrong. Please try again.')"
)
slide.placeholders[1].text = content

# Slide 5: Summary
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Summary"
content = (
    "• Real-time chat with AI\n"
    "• Responsive UI with loading indicators\n"
    "• Graceful handling of slow or failed responses\n"
    "• Modern, user-friendly interface"
)
slide.placeholders[1].text = content

# Slide 6: Notes for the User
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Notes for the User"
content = (
    "• Replace all screenshot placeholders with actual screenshots from your running app.\n"
    "• To simulate slow response, add `time.sleep(5)` in your backend before sending the AI response.\n"
    "• To simulate error, stop the backend or raise an exception in the chat route."
)
slide.placeholders[1].text = content

# Save the presentation
prs.save("Documents/AI_Chatbot_MVP_Demo.pptx")
print("Presentation created: Documents/AI_Chatbot_MVP_Demo.pptx") 